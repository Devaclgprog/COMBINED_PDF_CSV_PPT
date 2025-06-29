import os
import time
import tempfile
import streamlit as st
from datetime import datetime
from dotenv import load_dotenv
from PyPDF2 import PdfReader
import pdfplumber
import pandas as pd
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
import random

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

MODEL_NAME = "gemini-2.0-flash"
MAX_PDF_SIZE_MB = 50
MAX_SLIDES = 10
PROCESSING_CHUNK_SIZE = 15000
MIN_CONTENT_LENGTH = 100
PPT_FONT = "Calibri"

# ==================== NAVIGATION =====================
st.set_page_config(page_title="AI Utilities", page_icon="🤖")
st.title("📚 AI Powered Data Tools")

page = st.sidebar.radio("Choose a tool:", ["PDF Chat", "CSV Chat", "PPT From PDF"])

# ==================== PDF CHAT =====================
def pdf_chat():
    st.header("📄 Enhanced PDF Chat Assistant")
    uploaded_file = st.file_uploader("Upload your PDF file", type=["pdf"])

    if "pdf_chat" not in st.session_state:
        st.session_state.pdf_chat = {
            "messages": [],
            "file_processed": False,
            "gemini_file": None,
            "chat_session": None
        }

    if uploaded_file and not st.session_state.pdf_chat["file_processed"]:
        with st.spinner("Processing your PDF..."):
            try:
                with open("temp_upload.pdf", "wb") as f:
                    f.write(uploaded_file.getbuffer())
                file = genai.upload_file("temp_upload.pdf", mime_type="application/pdf")
                while file.state.name == "PROCESSING":
                    time.sleep(2)
                    file = genai.get_file(file.name)
                chat = genai.GenerativeModel(model_name=MODEL_NAME).start_chat(history=[{
                    "role": "user",
                    "parts": [file, "You are a PDF analysis assistant..."]
                }])
                st.session_state.pdf_chat.update({
                    "file_processed": True,
                    "gemini_file": file,
                    "chat_session": chat,
                    "messages": [{"role": "assistant", "content": "Hi! Ready to help with your PDF."}]
                })
                st.rerun()
            except Exception as e:
                st.error(str(e))

    for msg in st.session_state.pdf_chat["messages"]:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    if prompt := st.chat_input("Ask about the PDF..."):
        st.session_state.pdf_chat["messages"].append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        with st.chat_message("assistant"):
            with st.spinner("Analyzing..."):
                try:
                    response = st.session_state.pdf_chat["chat_session"].send_message(prompt)
                    st.markdown(response.text)
                except Exception as e:
                    st.markdown(f"Error: {str(e)}")
        st.session_state.pdf_chat["messages"].append({"role": "assistant", "content": response.text})

# ==================== CSV CHAT =====================
def csv_chat():
    st.header("📊 CSV Data Analysis Assistant")
    uploaded_file = st.file_uploader("Upload your CSV file", type=["csv"])

    if "csv_chat" not in st.session_state:
        st.session_state.csv_chat = {
            "messages": [],
            "file_processed": False,
            "gemini_file": None,
            "chat_session": None,
            "df": None
        }

    if uploaded_file and not st.session_state.csv_chat["file_processed"]:
        with st.spinner("Processing your CSV file..."):
            try:
                df = pd.read_csv(uploaded_file)
                with tempfile.NamedTemporaryFile(delete=False, suffix=".csv") as tmp:
                    df.to_csv(tmp.name, index=False)
                    tmp_path = tmp.name
                file = genai.upload_file(tmp_path, mime_type="text/csv")
                while file.state.name == "PROCESSING":
                    time.sleep(2)
                    file = genai.get_file(file.name)
                chat = genai.GenerativeModel(model_name=MODEL_NAME).start_chat(history=[{
                    "role": "user",
                    "parts": [file, "You are a data assistant..."]
                }])
                st.session_state.csv_chat.update({
                    "file_processed": True,
                    "gemini_file": file,
                    "chat_session": chat,
                    "df": df,
                    "messages": [{"role": "assistant", "content": "Hi! I have loaded your CSV file."}]
                })
                st.rerun()
            except Exception as e:
                st.error(str(e))

    if st.session_state.csv_chat["df"] is not None:
        if st.checkbox("Show Data Preview"):
            st.dataframe(st.session_state.csv_chat["df"].head())

    for msg in st.session_state.csv_chat["messages"]:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    if prompt := st.chat_input("Ask about your data..."):
        st.session_state.csv_chat["messages"].append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        with st.chat_message("assistant"):
            with st.spinner("Analyzing..."):
                try:
                    response = st.session_state.csv_chat["chat_session"].send_message(prompt)
                    st.markdown(response.text)
                except Exception as e:
                    st.markdown(f"Error: {str(e)}")
        st.session_state.csv_chat["messages"].append({"role": "assistant", "content": response.text})

# ==================== PPT FROM PDF =====================
def ppt_from_pdf():
    st.header("📊 PDF to PowerPoint Converter")
    model = genai.GenerativeModel(model_name=MODEL_NAME)
    pdf_file = st.file_uploader("Upload a PDF", type=["pdf"])

    if pdf_file:
        with st.spinner("Extracting text..."):
            text = ""
            try:
                with pdfplumber.open(pdf_file) as pdf:
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        if page_text:
                            text += f"\n\n[Page {i+1}]\n{page_text}"
                        if len(text) > PROCESSING_CHUNK_SIZE:
                            break
            except:
                reader = PdfReader(pdf_file)
                for page in reader.pages:
                    text += page.extract_text()
                    if len(text) > PROCESSING_CHUNK_SIZE:
                        break
            if not text:
                st.error("Failed to extract text.")
                return

        ppt_title = st.text_input("Presentation Title", "Business Report")

        if st.button("Generate Slides"):
            with st.spinner("Generating slide structure..."):
                prompt = f"""
You are an expert presentation designer. Based on the content below, create a PowerPoint structure for the title: '{ppt_title}'.
PDF CONTENT:
{text[:PROCESSING_CHUNK_SIZE]}
Return exactly 5 slides in this format:
**Slide 1: [Title Slide]**
* **Title:** "{ppt_title}"
* **Subtitle:** "[1-line summary]"
**Slide 2: [Introduction]**
* **Title:** "Intro title"
* **Bullet Points:**
    * Bullet 1
    * Bullet 2
...
"""
                slide_structure = model.generate_content(prompt).text
                st.code(slide_structure)

            if st.button("Generate PowerPoint"):
                with st.spinner("Creating presentation..."):
                    prs = Presentation()
                    prs.slide_width = Inches(13.333)
                    prs.slide_height = Inches(7.5)
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    slide.shapes.title.text = ppt_title
                    slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%d %B %Y')}"

                    slides = re.split(r'\*\*Slide \d+:.*?\*\*', slide_structure)
                    headers = re.findall(r'\*\*Slide \d+:.*?\*\*', slide_structure)
                    for i, header in enumerate(headers):
                        content = slides[i+1] if i+1 < len(slides) else ""
                        title_match = re.search(r'\*\*Title:\*\*\s*(.*)', content)
                        title = title_match.group(1).strip() if title_match else f"Slide {i+2}"
                        bullets = re.findall(r'\*\s+(.*)', content)
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = title
                        text_frame = slide.placeholders[1].text_frame
                        text_frame.clear()
                        for bullet in bullets:
                            para = text_frame.add_paragraph()
                            para.text = bullet
                            para.font.size = Pt(18)
                            para.font.color.rgb = RGBColor(0, 0, 0)
                            para.font.name = PPT_FONT
                            para.alignment = PP_ALIGN.LEFT

                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                        prs.save(tmp.name)
                        with open(tmp.name, "rb") as f:
                            ppt_bytes = f.read()
                        st.download_button(
                            label="📥 Download PowerPoint",
                            data=ppt_bytes,
                            file_name=f"{ppt_title.replace(' ', '_')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )

# ==================== MAIN =====================
if page == "PDF Chat":
    pdf_chat()
elif page == "CSV Chat":
    csv_chat()
elif page == "PPT From PDF":
    ppt_from_pdf()
